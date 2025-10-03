#!/usr/bin/env python3
"""
Script to create a sample PowerPoint file for testing the AI Document Chat app.
Run this script to generate a sample.pptx file.
"""

from pptx import Presentation
from pptx.util import Inches

def create_sample_presentation():
    """Create a sample PowerPoint presentation for testing"""
    
    # Create a new presentation
    prs = Presentation()
    
    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "AI and Machine Learning Overview"
    subtitle1.text = "A comprehensive introduction to artificial intelligence"
    
    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide layout
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "What is Artificial Intelligence?"
    content2.text = """• AI is the simulation of human intelligence in machines
• Machines are programmed to think and learn like humans
• Key areas include machine learning, natural language processing, and computer vision
• AI can perform tasks that typically require human intelligence"""
    
    # Slide 3: Content slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Types of Machine Learning"
    content3.text = """• Supervised Learning: Learning with labeled data
• Unsupervised Learning: Finding patterns without labels
• Reinforcement Learning: Learning through trial and error
• Deep Learning: Using neural networks with multiple layers"""
    
    # Slide 4: Content slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Real-World Applications"
    content4.text = """• Healthcare: Medical diagnosis and drug discovery
• Finance: Fraud detection and algorithmic trading
• Transportation: Autonomous vehicles
• Technology: Virtual assistants and recommendation systems
• Education: Personalized learning platforms"""
    
    # Slide 5: Conclusion
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "The Future of AI"
    content5.text = """• AI will continue to transform industries
• Ethical considerations are important
• Human-AI collaboration is key
• Continuous learning and adaptation required"""
    
    # Save the presentation
    prs.save('sample_presentation.pptx')
    print("✅ Sample PowerPoint presentation created: sample_presentation.pptx")

if __name__ == "__main__":
    create_sample_presentation()
