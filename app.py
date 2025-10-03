import streamlit as st
import openai
import os
from dotenv import load_dotenv
import tempfile
import json
from pptx import Presentation

# Load environment variables
load_dotenv()

# Configure Streamlit page
st.set_page_config(
    page_title="AI Document Chat",
    page_icon="📄",
    layout="wide"
)

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "document_content" not in st.session_state:
    st.session_state.document_content = ""
if "api_key" not in st.session_state:
    st.session_state.api_key = ""

def get_openai_client(api_key):
    """Initialize OpenAI client with the provided API key"""
    if not api_key:
        return None
    try:
        client = openai.OpenAI(api_key=api_key)
        # Test the API key by making a simple request
        client.models.list()
        return client
    except Exception as e:
        st.error(f"Error connecting to OpenAI API: {str(e)}")
        return None

def process_document(file_content, filename):
    """Process the uploaded document and extract text content"""
    try:
        if filename.endswith('.txt') or filename.endswith('.md'):
            # Handle text and markdown files
            return file_content.decode('utf-8')
        elif filename.endswith('.pptx'):
            # Handle PowerPoint files
            return process_pptx_file(file_content)
        else:
            # For other file types, try to decode as text
            return file_content.decode('utf-8')
    except Exception as e:
        st.error(f"Error processing document: {str(e)}")
        return ""

def process_pptx_file(file_content):
    """Extract comprehensive text content from a PowerPoint file with enhanced extraction"""
    try:
        # Create a temporary file to save the PowerPoint content
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(file_content)
            temp_file_path = temp_file.name
        
        # Load the presentation
        presentation = Presentation(temp_file_path)
        
        # Extract comprehensive content from all slides
        extracted_content = []
        total_slides = len(presentation.slides)
        
        for i, slide in enumerate(presentation.slides, 1):
            slide_content = extract_slide_content(slide, i)
            if slide_content:
                extracted_content.append(slide_content)
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        if not extracted_content:
            return "No text content found in the PowerPoint file."
        
        # Combine all slide content
        full_content = "\n\n".join(extracted_content)
        
        # Add presentation metadata
        metadata = f"PRESENTATION OVERVIEW:\nTotal Slides: {total_slides}\nContent Extracted: {len(extracted_content)} slides\n\n"
        
        return metadata + full_content
        
    except Exception as e:
        st.error(f"Error processing PowerPoint file: {str(e)}")
        return ""

def extract_slide_content(slide, slide_number):
    """Extract comprehensive content from a single slide"""
    slide_content = []
    slide_content.append(f"=== SLIDE {slide_number} ===")
    
    # Extract slide notes if available
    if hasattr(slide, 'notes_slide') and slide.notes_slide:
        notes_text = extract_text_from_shapes(slide.notes_slide.shapes)
        if notes_text.strip():
            slide_content.append(f"NOTES: {notes_text.strip()}")
    
    # Extract content from all shapes with detailed processing
    shape_content = []
    shape_count = 0
    
    for shape in slide.shapes:
        shape_count += 1
        shape_text = extract_text_from_shape(shape)
        
        # Debug: Log shape information
        shape_info = f"Shape {shape_count}: {type(shape).__name__}"
        if hasattr(shape, 'name'):
            shape_info += f" (name: {shape.name})"
        if hasattr(shape, 'shape_type'):
            shape_info += f" (type: {shape.shape_type})"
        
        if shape_text.strip():
            # Determine content type based on shape properties
            content_type = get_content_type(shape)
            if content_type:
                shape_content.append(f"{content_type}: {shape_text.strip()}")
            else:
                shape_content.append(shape_text.strip())
        else:
            # Log shapes with no text for debugging
            shape_content.append(f"DEBUG: {shape_info} - No text extracted")
    
    if shape_content:
        slide_content.extend(shape_content)
    
    # Only return slide content if it has meaningful text
    if len(slide_content) > 1:  # More than just the slide header
        return "\n".join(slide_content)
    
    return None

def extract_text_from_shape(shape):
    """Extract text from a single shape with comprehensive handling"""
    text_content = []
    
    try:
        # Handle different shape types - be more aggressive about text extraction
        
        # First, try to get any text from the shape
        if hasattr(shape, 'text') and shape.text and shape.text.strip():
            text_content.append(shape.text.strip())
        
        # Handle text frames (most common for detailed content)
        if hasattr(shape, 'text_frame') and shape.text_frame:
            frame_text = extract_text_frame_content(shape.text_frame)
            if frame_text and frame_text.strip():
                text_content.append(frame_text.strip())
        
        # Handle tables
        if hasattr(shape, 'table') and shape.table:
            table_text = extract_table_content(shape.table)
            if table_text and table_text.strip():
                text_content.append(table_text.strip())
        
        # Handle grouped shapes
        if hasattr(shape, 'shapes') and shape.shapes:
            group_text = extract_text_from_shapes(shape.shapes)
            if group_text and group_text.strip():
                text_content.append(group_text.strip())
        
        # Additional text extraction methods
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            try:
                if shape.text_frame and shape.text_frame.text:
                    text_content.append(shape.text_frame.text.strip())
            except:
                pass
        
        # Try to get text from any text-related attributes
        for attr in ['text', 'content', 'value']:
            if hasattr(shape, attr):
                try:
                    attr_value = getattr(shape, attr)
                    if attr_value and str(attr_value).strip():
                        text_content.append(str(attr_value).strip())
                except:
                    pass
    
    except Exception as e:
        # Continue processing other shapes if one fails
        pass
    
    # Join all text content and clean up
    result = " ".join(text_content)
    return result.strip() if result.strip() else ""

def extract_text_from_shapes(shapes):
    """Extract text from a collection of shapes"""
    text_parts = []
    for shape in shapes:
        shape_text = extract_text_from_shape(shape)
        if shape_text.strip():
            text_parts.append(shape_text.strip())
    return " ".join(text_parts)

def extract_table_content(table):
    """Extract content from a table"""
    table_content = []
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_text = cell.text.strip()
            if cell_text:
                row_content.append(cell_text)
        if row_content:
            table_content.append(" | ".join(row_content))
    
    if table_content:
        return "TABLE:\n" + "\n".join(table_content)
    return ""

def extract_text_frame_content(text_frame):
    """Extract content from a text frame with paragraph structure"""
    frame_content = []
    
    try:
        # Extract text from all paragraphs
        for paragraph in text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                # Check if this is a bullet point or special formatting
                if hasattr(paragraph, 'level') and paragraph.level > 0:
                    indent = "  " * paragraph.level
                    frame_content.append(f"{indent}• {para_text}")
                else:
                    frame_content.append(para_text)
        
        # Also try to get text from runs within paragraphs
        for paragraph in text_frame.paragraphs:
            if hasattr(paragraph, 'runs'):
                for run in paragraph.runs:
                    run_text = run.text.strip()
                    if run_text and run_text not in frame_content:
                        frame_content.append(run_text)
        
        # Try alternative text extraction
        if hasattr(text_frame, 'text') and text_frame.text:
            full_text = text_frame.text.strip()
            if full_text and full_text not in frame_content:
                frame_content.append(full_text)
    
    except Exception as e:
        # If detailed extraction fails, try basic text extraction
        try:
            if hasattr(text_frame, 'text') and text_frame.text:
                frame_content.append(text_frame.text.strip())
        except:
            pass
    
    return "\n".join(frame_content)

def get_content_type(shape):
    """Determine the type of content based on shape properties"""
    try:
        # Check if it's a title shape
        if hasattr(shape, 'placeholder_format'):
            if shape.placeholder_format.type == 1:  # Title placeholder
                return "TITLE"
            elif shape.placeholder_format.type == 2:  # Content placeholder
                return "CONTENT"
        
        # Check shape name for hints
        if hasattr(shape, 'name'):
            name_lower = shape.name.lower()
            if 'title' in name_lower:
                return "TITLE"
            elif 'content' in name_lower or 'body' in name_lower:
                return "CONTENT"
        
        # Check if it's a table
        if hasattr(shape, 'table'):
            return "TABLE"
        
        # Check if it's a text box
        if hasattr(shape, 'text_frame') and not hasattr(shape, 'table'):
            return "TEXT"
        
    except Exception:
        pass
    
    return None

def estimate_tokens(text):
    """Rough estimation of tokens (1 token ≈ 4 characters)"""
    return len(text) // 4

def create_slide_summaries(document_content):
    """Create summaries for each slide to reduce token usage"""
    slides = document_content.split("=== SLIDE")
    summaries = []
    
    for i, slide_content in enumerate(slides[1:], 1):  # Skip first empty split
        # Extract key information from each slide
        lines = slide_content.strip().split('\n')
        title = ""
        key_points = []
        
        for line in lines:
            if line.startswith("TITLE:"):
                title = line.replace("TITLE:", "").strip()
            elif line.startswith("CONTENT:") or line.startswith("TEXT:"):
                content = line.split(":", 1)[1].strip()
                if len(content) > 100:
                    content = content[:100] + "..."
                key_points.append(content)
            elif line.startswith("TABLE:"):
                key_points.append("Contains table data")
        
        # Create slide summary
        summary = f"Slide {i}: {title}" if title else f"Slide {i}"
        if key_points:
            summary += f" - Key points: {'; '.join(key_points[:3])}"  # Limit to 3 key points
        
        summaries.append(summary)
    
    return summaries

def get_relevant_slides(document_content, user_message):
    """Extract relevant slides based on user query"""
    import re
    
    # Look for slide number mentions in the user message
    slide_numbers = re.findall(r'slide\s+(\d+)', user_message.lower())
    
    if slide_numbers:
        # Extract specific slides mentioned
        slides = document_content.split("=== SLIDE")
        relevant_content = []
        
        for slide_num in slide_numbers:
            slide_index = int(slide_num)
            # Adjust for 0-based indexing and ensure we have the slide
            if slide_index > 0 and slide_index < len(slides):
                slide_content = slides[slide_index]
                relevant_content.append(f"=== SLIDE{slide_content}")
        
        if relevant_content:
            return "\n\n".join(relevant_content)
    
    # If no specific slides mentioned, return first few slides and summary
    slides = document_content.split("=== SLIDE")
    if len(slides) > 1:
        # Return first 5 slides + summary of the rest
        first_slides = slides[1:6]  # First 5 slides
        remaining_slides = slides[6:] if len(slides) > 6 else []
        
        content_parts = [f"=== SLIDE{slide}" for slide in first_slides]
        
        if remaining_slides:
            summaries = create_slide_summaries(document_content)
            content_parts.append(f"\nSUMMARY OF REMAINING SLIDES:\n" + "\n".join(summaries[5:]))
        
        return "\n\n".join(content_parts)
    
    return document_content

def chat_with_ai(client, user_message, document_context=""):
    """Send a message to OpenAI with smart content management for large presentations"""
    try:
        system_prompt = "You are a helpful AI assistant specialized in analyzing presentations and documents. "
        
        if document_context:
            # Check token count and implement smart chunking
            estimated_tokens = estimate_tokens(document_context)
            
            if estimated_tokens > 12000:  # Conservative limit for gpt-3.5-turbo
                # Use smart content selection
                relevant_content = get_relevant_slides(document_context, user_message)
                
                system_prompt += "The user has uploaded a large presentation. "
                system_prompt += "You have access to relevant slide content based on their question. "
                system_prompt += "If they ask about specific slides, provide detailed information from those slides. "
                system_prompt += "If they ask general questions, provide a comprehensive overview based on the available content.\n\n"
                system_prompt += f"RELEVANT PRESENTATION CONTENT:\n{relevant_content}\n\n"
                
                # Add note about full presentation
                system_prompt += "NOTE: This is a large presentation. If you need information from other slides, ask the user to specify which slides they're interested in.\n\n"
                
                # Add debugging info for development
                if "slide" in user_message.lower():
                    system_prompt += "IMPORTANT: The user is asking about a specific slide. Make sure to provide comprehensive details from the slide content above. Look for all text, tables, and content within that slide.\n\n"
            else:
                system_prompt += f"Here is the content from the user's uploaded document:\n{document_context}\n\n"
            
            system_prompt += "Please provide detailed, accurate responses about the document content. "
            system_prompt += "When referencing specific slides, mention the slide number. "
            system_prompt += "If asked about specific details, provide comprehensive information from the relevant slides."
        else:
            system_prompt += "Please help the user with their questions."
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            max_tokens=1500,
            temperature=0.7
        )
        
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"

# Main app layout
st.title("📄 AI Document Chat")
st.markdown("Upload a document and chat with AI about its content!")

# Create two columns
col1, col2 = st.columns([1, 2])

# Left column - Controls
with col1:
    st.header("Settings")
    
    # API Key input
    st.subheader("OpenAI API Key")
    api_key_input = st.text_input(
        "Enter your OpenAI API key:",
        value=st.session_state.api_key,
        type="password",
        help="You can also set the OPENAI_API_KEY environment variable"
    )
    
    # Update session state
    if api_key_input:
        st.session_state.api_key = api_key_input
    
    # Check for environment variable
    env_api_key = os.getenv("OPENAI_API_KEY")
    if env_api_key and not api_key_input:
        st.session_state.api_key = env_api_key
        st.info("Using API key from environment variable")
    
    # File upload
    st.subheader("Upload Document")
    uploaded_file = st.file_uploader(
        "Choose a document to upload:",
        type=['txt', 'md', 'pptx'],
        help="Supports .txt, .md, and .pptx files"
    )
    
    if uploaded_file is not None:
        # Show processing indicator for large files
        file_size_mb = len(uploaded_file.read()) / (1024 * 1024)
        uploaded_file.seek(0)  # Reset file pointer
        
        if file_size_mb > 10:  # Large file
            st.info(f"📊 Processing large file ({file_size_mb:.1f} MB). This may take a moment...")
        
        # Read the file content
        file_content = uploaded_file.read()
        
        # Show progress for PowerPoint files
        if uploaded_file.name.endswith('.pptx'):
            with st.spinner("🔍 Extracting content from PowerPoint slides..."):
                document_text = process_document(file_content, uploaded_file.name)
        else:
            document_text = process_document(file_content, uploaded_file.name)
        
        if document_text:
            st.session_state.document_content = document_text
            
            # Show success with file stats
            content_length = len(document_text)
            st.success(f"✅ Document '{uploaded_file.name}' processed successfully!")
            
            # Show document statistics
            if uploaded_file.name.endswith('.pptx'):
                # Extract slide count from the processed content
                slide_count = document_text.count("=== SLIDE")
                st.info(f"📊 Extracted content from {slide_count} slides ({content_length:,} characters)")
                
                # Show helpful tips for large presentations
                if content_length > 50000:  # Large presentation
                    st.warning("⚠️ **Large Presentation Detected**")
                    st.markdown("""
                    **Tips for analyzing large presentations:**
                    - Ask about specific slides: *"What's on slide 15?"*
                    - Request summaries: *"Summarize slides 10-20"*
                    - Focus on sections: *"What are the key points from the first 10 slides?"*
                    """)
            
            # Show preview with better formatting
            preview_text = document_text[:1000] + "..." if len(document_text) > 1000 else document_text
            st.text_area(
                "Document Preview:",
                value=preview_text,
                height=200,
                disabled=True,
                help="This is a preview of the extracted content. The AI has access to the full document."
            )
        else:
            st.error("❌ Failed to process the document")
    
    # Clear document button
    if st.session_state.document_content:
        if st.button("🗑️ Clear Document"):
            st.session_state.document_content = ""
            st.session_state.messages = []
            st.rerun()

# Right column - Chat interface
with col2:
    st.header("AI Chat")
    
    # Check if API key is available
    if not st.session_state.api_key:
        st.warning("⚠️ Please enter your OpenAI API key in the left panel to start chatting.")
    else:
        # Initialize OpenAI client
        client = get_openai_client(st.session_state.api_key)
        
        if client is None:
            st.error("❌ Unable to connect to OpenAI API. Please check your API key.")
        else:
            # Show helpful message for large presentations
            if st.session_state.document_content and len(st.session_state.document_content) > 50000:
                st.info("💡 **Large presentation loaded!** For best results, ask about specific slides or sections (e.g., 'What's on slide 15?' or 'Summarize the first 10 slides').")
            
            # Display chat messages
            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            
            # Chat input
            if prompt := st.chat_input("Ask me anything about your document or any other topic!"):
                # Add user message to chat history
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)
                
                # Get AI response
                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        response = chat_with_ai(
                            client, 
                            prompt, 
                            st.session_state.document_content
                        )
                    st.markdown(response)
                
                # Add assistant response to chat history
                st.session_state.messages.append({"role": "assistant", "content": response})

# Footer
st.markdown("---")
st.markdown("**Instructions:**")
st.markdown("1. Enter your OpenAI API key in the left panel")
st.markdown("2. Upload a document (supports .txt, .md, and .pptx files)")
st.markdown("3. Start chatting with the AI about your document!")
